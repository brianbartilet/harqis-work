"""
workflows/hfl/tasks/time_capsule.py

Time-capsule synthesizer — the COLLECT half of an on-demand, time-ranged
filesystem ingest into the Homework-for-Life corpus.

Given a directory (default the harqis-server data root ``/Volumes/harqis-data``)
and a flexible date *period* (``May 2020``, ``August 1, 2002``,
``June-July 2019``, ``2019-06-01..2019-07-31``, ``since 2020-01-01``), this
module sweeps the tree (and every subdirectory) for files whose mtime falls in
the window and extracts a compact, bounded representation of each:

  - text / logs / code / config → a head snippet read directly,
  - documents (pdf/docx/xlsx/pptx) → extracted text (optional parsers; degrade
    to metadata-only when a parser is missing),
  - images / video → a short Haiku vision caption (reuses analyze_media's
    bounded encoders; degrades to metadata-only without cv2 / Anthropic),
  - audio → metadata-only (no transcriber wired yet — see README),
  - anything else → metadata only.

It writes two artifacts and prints the digest to stdout:
  - ``<repo>/logs/time-capsule/<slug>.manifest.json`` — the full record,
  - ``<repo>/logs/time-capsule/<slug>.digest.md``    — a compact, LLM-friendly
    digest (also mirrored into ``<corpus>/.time-capsule/`` for provenance).

This is the HYBRID design's first phase: the task COLLECTS + EXTRACTS; the
``/time-capsule-synthesizer`` skill (Claude) reads the digest and SYNTHESIZES
one period rollup ``HflEntry``, then calls :func:`run_write` here to dual-write
it (corpus + ES) via the shared ``capture.append_entry`` helper.

Cost: per-file captions use Haiku only (``claude-haiku-4-5-20251001``) and are
bounded by ``max_caption_files`` + ``max_files``. Never raise the Anthropic
DEFAULT_MODEL. One bad file never aborts the sweep (logged, metadata-only).

References:
- workflows/hfl/tasks/analyze_media.py — the windowed vision sweep this mirrors
  (its ``_encode_image`` / ``_video_blocks`` encoders are reused directly).
- workflows/dumps/files.py — ``iter_recent_files`` (the mtime-windowed walk).
- workflows/hfl/tasks/capture.py — ``_build_entry`` / ``append_entry`` /
  ``resolve_corpus_dir`` (the dual-write the rollup is written through).
"""

from __future__ import annotations

import calendar
import json
import re
import sys
from collections import Counter, defaultdict
from dataclasses import dataclass
from datetime import datetime, timedelta
from pathlib import Path
from typing import Any, Optional

from core.apps.sprout.app.celery import SPROUT
from core.apps.es_logging.app.elasticsearch import log_result
from core.utilities.logging.custom_logger import create_logger

from apps.antropic.config import get_config as get_anthropic_config
from apps.antropic.references.web.base_api_service import BaseApiServiceAnthropic

from workflows.dumps.files import iter_recent_files
from workflows.hfl.prompts import load_prompt
from workflows.hfl.tasks.capture import (
    _build_entry,
    append_entry,
    resolve_corpus_dir,
)
# Reuse the proven, size-bounded vision encoders + extension sets rather than
# duplicating them — analyze_media already handles cv2-optional degradation.
from workflows.hfl.tasks.analyze_media import (
    _encode_image,
    _video_blocks,
    _IMAGE_EXTS,
    _VIDEO_EXTS,
)

_log = create_logger("hfl.time_capsule")

# __file__ = repo/workflows/hfl/tasks/time_capsule.py → parents[3] = repo root.
_REPO_ROOT = Path(__file__).resolve().parents[3]
_OUT_DIR = _REPO_ROOT / "logs" / "time-capsule"

_DEFAULT_HAIKU = "claude-haiku-4-5-20251001"
# Default root: the harqis-server data volume (see machines.local.toml [dumps]
# / [harqis-server.env_vars] HFL_CORPUS_PATH, both under /Volumes/harqis-data).
DEFAULT_ROOT = "/Volumes/harqis-data"

# File-kind extension buckets (lower-case, dot-prefixed). Images/video are
# imported from analyze_media so the two stay in lockstep.
_TEXT_EXTS = {
    ".txt", ".log", ".md", ".markdown", ".rst", ".csv", ".tsv", ".json",
    ".yaml", ".yml", ".xml", ".ini", ".cfg", ".conf", ".toml", ".html",
    ".htm", ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".c", ".cpp",
    ".h", ".hpp", ".cs", ".go", ".rb", ".rs", ".php", ".sh", ".ps1", ".bat",
    ".sql", ".env", ".gitignore", ".srt", ".vtt",
}
_DOC_EXTS = {".pdf", ".docx", ".doc", ".xlsx", ".xls", ".pptx", ".ppt", ".odt"}
_AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".aac", ".flac", ".ogg", ".opus", ".wma"}

_CAPTION_PROMPT = load_prompt("time_capsule_caption").strip()


# ── Period parsing ────────────────────────────────────────────────────────────

# month name / abbrev → number (1-12). Includes the common "sept" alias.
_MONTHS: dict[str, int] = {}
for _i in range(1, 13):
    _MONTHS[calendar.month_name[_i].lower()] = _i
    _MONTHS[calendar.month_abbr[_i].lower()] = _i
_MONTHS["sept"] = 9


@dataclass(frozen=True)
class Window:
    """A half-open ``[start, end)`` time window plus a filename-safe label."""
    start: datetime
    end: datetime
    label: str


def _first_of_next_month(year: int, month: int) -> datetime:
    return datetime(year + 1, 1, 1) if month == 12 else datetime(year, month + 1, 1)


def _resolve_token(token: str, *, default_year: Optional[int] = None) -> tuple[datetime, datetime]:
    """Resolve one date token to a half-open ``[start, end)`` pair.

    Granularity is inferred from the token: a bare year spans the year, a
    ``YYYY-MM`` or ``Month YYYY`` spans the month, a full date spans one day.
    ``default_year`` lets a bare month name (used inside ``Mon-Mon YYYY``)
    inherit the range's year.
    """
    tok = token.strip().strip(",").strip()
    low = tok.lower()

    m = re.fullmatch(r"(\d{4})[-/.](\d{1,2})[-/.](\d{1,2})", low)
    if m:
        y, mo, d = (int(x) for x in m.groups())
        s = datetime(y, mo, d)
        return s, s + timedelta(days=1)

    m = re.fullmatch(r"(\d{4})[-/](\d{1,2})", low)
    if m:
        y, mo = int(m.group(1)), int(m.group(2))
        return datetime(y, mo, 1), _first_of_next_month(y, mo)

    m = re.fullmatch(r"(\d{4})", low)
    if m:
        y = int(m.group(1))
        return datetime(y, 1, 1), datetime(y + 1, 1, 1)

    # Month D, YYYY  /  Month D YYYY
    m = re.fullmatch(r"([a-z]+)\.?\s+(\d{1,2})(?:st|nd|rd|th)?,?\s+(\d{4})", low)
    if m and m.group(1) in _MONTHS:
        s = datetime(int(m.group(3)), _MONTHS[m.group(1)], int(m.group(2)))
        return s, s + timedelta(days=1)

    # D Month YYYY  (e.g. "1 August 2002")
    m = re.fullmatch(r"(\d{1,2})(?:st|nd|rd|th)?\s+([a-z]+)\.?\s+(\d{4})", low)
    if m and m.group(2) in _MONTHS:
        s = datetime(int(m.group(3)), _MONTHS[m.group(2)], int(m.group(1)))
        return s, s + timedelta(days=1)

    # Month YYYY
    m = re.fullmatch(r"([a-z]+)\.?\s+(\d{4})", low)
    if m and m.group(1) in _MONTHS:
        y, mo = int(m.group(2)), _MONTHS[m.group(1)]
        return datetime(y, mo, 1), _first_of_next_month(y, mo)

    # bare Month, inheriting a range year (e.g. "june" in "june-july 2019")
    if low in _MONTHS and default_year is not None:
        mo = _MONTHS[low]
        return datetime(default_year, mo, 1), _first_of_next_month(default_year, mo)

    raise ValueError(f"unrecognized date token: {token!r}")


def parse_period(text: str, *, now: Optional[datetime] = None) -> Window:
    """Parse a flexible period string into a :class:`Window` (``[start, end)``).

    Supported forms (case-insensitive):
      - ``May 2020`` / ``2020-05`` / ``2020`` (whole month / year)
      - ``August 1, 2002`` / ``1 Aug 2002`` / ``2002-08-01`` (single day)
      - ``June-July 2019`` (month range sharing a year)
      - ``2019-06-01..2019-07-31`` / ``Jan 2020 to Mar 2020`` (explicit range)
      - ``since 2020-01-01`` / ``from May 2020`` / ``2020-01-01 to today``
      - ``-30d`` / ``last 30 days`` (relative to now)
    """
    now = now or datetime.now()
    raw = (text or "").strip()
    low = raw.lower().strip()
    if not low:
        raise ValueError("empty period — pass e.g. 'May 2020' or 'since 2020-01-01'")

    m = re.fullmatch(r"-?(\d+)\s*d(?:ays)?", low) or re.fullmatch(r"last\s+(\d+)\s+days?", low)
    if m:
        n = int(m.group(1))
        return Window(now - timedelta(days=n), now, f"last-{n}d")

    m = re.match(r"(?:since|from|after)\s+(.+)$", low)
    if m:
        s, _ = _resolve_token(m.group(1))
        return Window(s, now, f"{s:%Y-%m-%d}_to_today")

    m = re.match(r"(.+?)\s+(?:to|until|up to|through|thru)\s+(?:today|now|present)$", low)
    if m:
        s, _ = _resolve_token(m.group(1))
        return Window(s, now, f"{s:%Y-%m-%d}_to_today")

    # Mon-Mon YYYY (the year applies to both months).
    m = re.fullmatch(r"([a-z]+)\s*[-–—]\s*([a-z]+)\s+(\d{4})", low)
    if m and m.group(1) in _MONTHS and m.group(2) in _MONTHS:
        y = int(m.group(3))
        s, _ = _resolve_token(m.group(1), default_year=y)
        _, e = _resolve_token(m.group(2), default_year=y)
        if e <= s:  # tolerate reversed ("Dec-Jan") by swapping
            s, e = _resolve_token(m.group(2), default_year=y)[0], _resolve_token(m.group(1), default_year=y)[1]
        return Window(s, e, f"{m.group(1)[:3]}-{m.group(2)[:3]}-{y}")

    # Explicit two-sided range. ".." first (ISO-safe); then spaced separators
    # so a hyphen inside an ISO date isn't mistaken for the range separator.
    for sep in ("..", " to ", " until ", " through ", " thru ", " – ", " — ", " - "):
        if sep in low:
            a, b = low.split(sep, 1)
            sa, _ = _resolve_token(a)
            _, eb = _resolve_token(b, default_year=sa.year)
            last = eb - timedelta(days=1)
            return Window(sa, eb, f"{sa:%Y%m%d}-{last:%Y%m%d}")

    s, e = _resolve_token(low)
    last = e - timedelta(days=1)
    label = f"{s:%Y%m%d}" if (e - s) <= timedelta(days=1) else f"{s:%Y%m%d}-{last:%Y%m%d}"
    return Window(s, e, label)


def _slug(label: str) -> str:
    """Filename-safe slug for a window label."""
    return re.sub(r"[^a-z0-9]+", "-", (label or "").lower()).strip("-") or "period"


# ── Per-file extraction ─────────────────────────────────────────────────────

def _classify(suffix: str) -> str:
    s = (suffix or "").lower()
    if s in _IMAGE_EXTS:
        return "image"
    if s in _VIDEO_EXTS:
        return "video"
    if s in _AUDIO_EXTS:
        return "audio"
    if s in _DOC_EXTS:
        return "document"
    if s in _TEXT_EXTS:
        return "text"
    return "other"


def _safe_size(path: Path) -> int:
    try:
        return path.stat().st_size
    except OSError:
        return 0


def _read_text(path: Path, limit: int) -> Optional[str]:
    """Read a bounded head snippet of a text file (utf-8, errors replaced)."""
    try:
        with path.open("r", encoding="utf-8", errors="replace") as fh:
            data = fh.read(limit + 1)
    except OSError:
        return None
    data = data.strip()
    if not data:
        return None
    if len(data) > limit:
        data = data[:limit].rstrip() + " …[truncated]"
    return data


def _extract_document(path: Path, limit: int) -> Optional[str]:
    """Best-effort text extraction from a document. None when no parser /
    no text — the caller records metadata only (never crashes the sweep)."""
    suffix = path.suffix.lower()
    try:
        if suffix == ".pdf":
            from pypdf import PdfReader  # optional dep
            reader = PdfReader(str(path))
            parts = [(pg.extract_text() or "") for pg in reader.pages[:20]]
            text = "\n".join(parts)
        elif suffix == ".docx":
            import docx  # python-docx, optional dep
            doc = docx.Document(str(path))
            text = "\n".join(p.text for p in doc.paragraphs)
        elif suffix == ".xlsx":
            import openpyxl  # optional dep
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
            rows: list[str] = []
            for ws in wb.worksheets[:3]:
                for r in ws.iter_rows(max_row=200, values_only=True):
                    cells = [str(c) for c in r if c is not None]
                    if cells:
                        rows.append("\t".join(cells))
                    if len(rows) >= 200:
                        break
            wb.close()
            text = "\n".join(rows)
        elif suffix == ".pptx":
            from pptx import Presentation  # python-pptx, optional dep
            prs = Presentation(str(path))
            chunks: list[str] = []
            for slide in prs.slides[:40]:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip():
                        chunks.append(shape.text_frame.text.strip())
            text = "\n".join(chunks)
        else:
            return None  # .doc/.xls/.ppt/.odt — legacy formats not parsed in v1
    except ImportError:
        return None
    except Exception as exc:  # noqa: BLE001 - a corrupt doc is metadata-only, not fatal
        _log.debug("time_capsule: document extract failed for %s (%s)", path, exc)
        return None

    text = (text or "").strip()
    if not text:
        return None
    if len(text) > limit:
        text = text[:limit].rstrip() + " …[truncated]"
    return text


def _make_client(cfg_id: str) -> Optional[BaseApiServiceAnthropic]:
    """Build the Anthropic client, or None — captions degrade to metadata."""
    try:
        client = BaseApiServiceAnthropic(get_anthropic_config(cfg_id))
        return client if getattr(client, "base_client", None) else None
    except Exception as exc:  # noqa: BLE001 - vision is bonus, never fatal
        _log.warning("time_capsule: Anthropic unavailable — captions skipped (%s)", exc)
        return None


def _caption_media(path: Path, kind: str, mtime: datetime, rel: str, *,
                   client: BaseApiServiceAnthropic, model: str,
                   frames_per_video: int, max_tokens: int = 220) -> Optional[str]:
    """Short Haiku vision caption for one image / video. None when unreadable."""
    if kind == "image":
        block = _encode_image(path)
        blocks = [block] if block else []
    else:
        blocks = _video_blocks(path, frames_per_video)
    if not blocks:
        return None
    instruction = {
        "type": "text",
        "text": (
            f"File: {path.name}\n"
            f"Captured: {mtime.strftime('%Y-%m-%d %H:%M')}\n"
            f"Folder: {rel}\n\n"
            "Describe this artifact for a memory log per your instructions. "
            "Reply with the caption only."
        ),
    }
    resp = client.send_messages(
        messages=[{"role": "user", "content": [*blocks, instruction]}],
        model=model, max_tokens=max_tokens, system=_CAPTION_PROMPT,
    )
    text = resp.content[0].text if resp and resp.content else ""
    return text.strip() or None


# ── Collection ──────────────────────────────────────────────────────────────

def collect_archive(
    root: str | Path,
    window: Window,
    *,
    max_files: int = 500,
    max_caption_files: int = 80,
    frames_per_video: int = 3,
    caption_model: str = _DEFAULT_HAIKU,
    cfg_id: str = "ANTHROPIC",
    do_caption: bool = True,
    snippet_chars: int = 1200,
) -> dict[str, Any]:
    """Sweep ``root`` for files in ``window`` and return a bounded manifest.

    The manifest is the COLLECT phase output: per-file kind + extracted text /
    vision caption / metadata, plus per-day and per-kind aggregates. It never
    raises on a bad file — that file is recorded metadata-only with a ``note``.
    """
    root_path = Path(root).expanduser()
    collected = sorted(
        iter_recent_files([root_path], window.start, window.end),
        key=lambda c: c.mtime, reverse=True,
    )
    total_in_window = len(collected)
    selected = collected[:max_files]

    files: list[dict] = []
    by_kind: Counter = Counter()
    by_day: dict[str, dict] = defaultdict(lambda: {"files": 0, "kinds": Counter()})
    client: Optional[BaseApiServiceAnthropic] = None
    captioned = errored = 0

    for c in selected:
        kind = _classify(c.path.suffix)
        rec: dict[str, Any] = {
            "path": str(c.path),
            "name": c.path.name,
            "rel": c.relative.as_posix(),
            "mtime": c.mtime.strftime("%Y-%m-%d %H:%M"),
            "kind": kind,
            "bytes": _safe_size(c.path),
        }
        try:
            if kind == "text":
                rec["text"] = _read_text(c.path, snippet_chars)
            elif kind == "document":
                txt = _extract_document(c.path, snippet_chars)
                if txt:
                    rec["text"] = txt
                else:
                    rec["note"] = "no text extracted (parser missing or empty)"
            elif kind in ("image", "video"):
                if do_caption and captioned < max_caption_files:
                    if client is None:
                        client = _make_client(cfg_id)
                    if client is not None:
                        cap = _caption_media(
                            c.path, kind, c.mtime, c.relative.as_posix(),
                            client=client, model=caption_model,
                            frames_per_video=frames_per_video,
                        )
                        if cap:
                            rec["caption"] = cap
                            captioned += 1
                        else:
                            rec["note"] = "no caption (unreadable media or cv2 missing)"
                    else:
                        rec["note"] = "caption skipped (Anthropic unavailable)"
                else:
                    rec["note"] = "caption skipped (budget reached or disabled)"
            elif kind == "audio":
                rec["note"] = "audio not transcribed (no transcriber configured)"
            else:
                rec["note"] = "binary / unknown type — metadata only"
        except Exception as exc:  # noqa: BLE001 - one bad file never aborts the batch
            errored += 1
            rec["note"] = f"error: {exc}"[:200]
            _log.warning("time_capsule: failed on %s — %s", c.path, exc)

        files.append(rec)
        by_kind[kind] += 1
        day = c.mtime.strftime("%Y-%m-%d")
        by_day[day]["files"] += 1
        by_day[day]["kinds"][kind] += 1

    return {
        "root": str(root_path),
        "root_reachable": root_path.exists() and root_path.is_dir(),
        "window": {
            "start": window.start.isoformat(timespec="minutes"),
            "end": window.end.isoformat(timespec="minutes"),
            "label": window.label,
        },
        "generated_at": datetime.now().isoformat(timespec="seconds"),
        "counts": {
            "total_in_window": total_in_window,
            "analyzed": len(files),
            "captioned": captioned,
            "errored": errored,
            "by_kind": dict(by_kind),
        },
        "by_day": {
            d: {"files": v["files"], "kinds": dict(v["kinds"])}
            for d, v in sorted(by_day.items())
        },
        "files": files,
    }


# ── Digest rendering ────────────────────────────────────────────────────────

def render_digest(manifest: dict, *, max_file_lines: int = 200,
                  snippet_chars: int = 220) -> str:
    """Compact, LLM-friendly Markdown digest of a manifest (the synthesis input)."""
    w = manifest["window"]
    c = manifest["counts"]
    out: list[str] = []
    out.append(f"# Time Capsule — {w['label']}")
    out.append("")
    out.append(f"- Root: `{manifest['root']}`"
               + ("" if manifest.get("root_reachable") else "  ⚠️ NOT REACHABLE on this host"))
    out.append(f"- Window: {w['start']} → {w['end']} (end-exclusive)")
    out.append(f"- Files in window: {c['total_in_window']} · analyzed: {c['analyzed']} "
               f"· captioned: {c['captioned']} · errored: {c['errored']}")
    kinds = ", ".join(f"{k}={v}" for k, v in sorted(c["by_kind"].items())) or "(none)"
    out.append(f"- By kind: {kinds}")
    out.append("")

    if manifest["by_day"]:
        out.append("## By day")
        out.append("")
        for day, v in manifest["by_day"].items():
            kinds = ", ".join(f"{k}×{n}" for k, n in sorted(v["kinds"].items()))
            out.append(f"- {day}: {v['files']} file(s) — {kinds}")
        out.append("")

    out.append("## Files")
    out.append("")
    shown = manifest["files"][:max_file_lines]
    for f in shown:
        detail = f.get("text") or f.get("caption") or f.get("note") or ""
        detail = " ".join(detail.split())
        if len(detail) > snippet_chars:
            detail = detail[:snippet_chars].rstrip() + " …"
        out.append(f"- [{f['kind']}] {f['mtime']} `{f['rel']}` — {detail}")
    remaining = len(manifest["files"]) - len(shown)
    if remaining > 0:
        out.append(f"- …and {remaining} more analyzed file(s) (see manifest.json).")
    out.append("")
    return "\n".join(out) + "\n"


# ── Orchestration entry points (driven by the skill via the venv python) ──────

def _emit(text: str) -> None:
    """Print to stdout robustly across console encodings.

    The digest contains non-ASCII glyphs (→, …, ×). A Windows console defaults
    to cp1252, where ``print`` would raise UnicodeEncodeError, so write UTF-8
    bytes straight to the buffer. Under pytest capture (no ``.buffer``) fall
    back to a replace-encoded ``print``.
    """
    try:
        sys.stdout.buffer.write((text + "\n").encode("utf-8"))
        sys.stdout.buffer.flush()
    except (AttributeError, ValueError):
        try:
            print(text)
        except UnicodeEncodeError:
            enc = getattr(sys.stdout, "encoding", None) or "utf-8"
            print(text.encode(enc, "replace").decode(enc, "replace"))


def _write_artifacts(manifest: dict, digest: str, slug: str) -> dict[str, str]:
    """Persist manifest.json + digest.md locally and (best-effort) mirror the
    digest into ``<corpus>/.time-capsule/`` for provenance on the corpus host."""
    _OUT_DIR.mkdir(parents=True, exist_ok=True)
    manifest_path = _OUT_DIR / f"{slug}.manifest.json"
    digest_path = _OUT_DIR / f"{slug}.digest.md"
    manifest_path.write_text(json.dumps(manifest, indent=2, ensure_ascii=False), encoding="utf-8")
    digest_path.write_text(digest, encoding="utf-8")
    try:
        prov = resolve_corpus_dir() / ".time-capsule"
        prov.mkdir(parents=True, exist_ok=True)
        (prov / f"{slug}.digest.md").write_text(digest, encoding="utf-8")
    except Exception as exc:  # noqa: BLE001 - provenance copy is best-effort
        _log.debug("time_capsule: provenance copy skipped (%s)", exc)
    return {"manifest_path": str(manifest_path), "digest_path": str(digest_path)}


def run_collect(*, root: str = DEFAULT_ROOT, period: str,
                max_files: int = 500, max_caption_files: int = 80,
                frames_per_video: int = 3, do_caption: bool = True,
                caption_model: str = _DEFAULT_HAIKU, cfg_id: str = "ANTHROPIC",
                now: Optional[datetime] = None) -> dict[str, Any]:
    """COLLECT phase: parse the period, sweep the tree, write artifacts, and
    print the digest. Returns a small summary the skill reads from stdout.

    A suggested ``when_iso`` (the period's last day, clamped to today) is
    returned so the synthesized rollup lands chronologically in the corpus.
    """
    now = now or datetime.now()
    window = parse_period(period, now=now)
    slug = _slug(window.label)

    root_path = Path(root).expanduser()
    if not (root_path.exists() and root_path.is_dir()):
        return {
            "ok": False,
            "reason": "root-unreachable",
            "root": str(root_path),
            "message": (
                f"Directory not reachable on this host: {root_path}. Run the skill on "
                f"the machine that mounts it (harqis-server for /Volumes/harqis-data), "
                f"or pass --root <a path on this machine>."
            ),
        }

    manifest = collect_archive(
        root_path, window,
        max_files=max_files, max_caption_files=max_caption_files,
        frames_per_video=frames_per_video, do_caption=do_caption,
        caption_model=caption_model, cfg_id=cfg_id,
    )
    digest = render_digest(manifest)
    paths = _write_artifacts(manifest, digest, slug)

    # Suggested entry date: last day in the window, clamped to today.
    last_day = window.end - timedelta(days=1)
    when = min(last_day, now)
    when_iso = when.replace(microsecond=0).isoformat()

    _emit(digest)
    _emit(f"[time_capsule] manifest: {paths['manifest_path']}")
    _emit(f"[time_capsule] digest:   {paths['digest_path']}")
    _emit(f"[time_capsule] suggested_when_iso: {when_iso}")

    return {
        "ok": True,
        "slug": slug,
        "window": manifest["window"],
        "counts": manifest["counts"],
        "suggested_when_iso": when_iso,
        **paths,
    }


def run_write(*, synthesis_path: str, source: str = "time-capsule") -> dict[str, Any]:
    """WRITE phase: dual-write the synthesized rollup entry (corpus + ES).

    ``synthesis_path`` is a JSON file the skill produced with the keys:
    ``moment``, ``what_happened``, ``why_it_stayed``, ``possible_use``,
    ``tags`` (list), ``references`` (list of file paths), ``when_iso``.
    """
    data = json.loads(Path(synthesis_path).read_text(encoding="utf-8"))
    moment = (data.get("moment") or "").strip()
    if not moment:
        return {"ok": False, "reason": "empty-moment", "entries_written": 0}

    when_iso = data.get("when_iso")
    try:
        when = datetime.fromisoformat(when_iso) if when_iso else datetime.now()
    except ValueError:
        when = datetime.now()

    entry = _build_entry(
        when=when,
        moment=moment,
        what_happened=str(data.get("what_happened") or "").strip(),
        why_it_stayed=str(data.get("why_it_stayed") or "").strip(),
        possible_use=str(data.get("possible_use") or "").strip(),
        tags=[str(t).strip().lstrip("#") for t in (data.get("tags") or []) if str(t).strip()],
        references=[str(r).strip() for r in (data.get("references") or []) if str(r).strip()],
    )
    corpus_dir = resolve_corpus_dir()
    corpus_dir.mkdir(parents=True, exist_ok=True)
    day_file = corpus_dir / f"{when.strftime('%Y-%m-%d')}.md"
    bytes_written, doc_id = append_entry(day_file, entry, source=source, synthesized=True)

    _log.info("time_capsule: rollup written → %s (doc_id=%s)", day_file, doc_id)
    result = {
        "ok": True,
        "entries_written": 1,
        "path": str(day_file),
        "doc_id": doc_id,
        "indexed": doc_id is not None,
        "bytes_written": bytes_written,
        "references": len(entry.references),
    }
    _emit(f"[time_capsule] wrote entry -> {day_file}")
    _emit(f"[time_capsule] es_doc_id: {doc_id}")
    return result


# ── Celery task (dispatch path: run the sweep on the HFL queue / corpus host) ─

@SPROUT.task()
@log_result()
def collect_time_capsule(*, root: str = DEFAULT_ROOT, period: str,
                         max_files: int = 500, max_caption_files: int = 80,
                         frames_per_video: int = 3, do_caption: bool = True,
                         cfg_id__anthropic: str = "ANTHROPIC",
                         model: str = _DEFAULT_HAIKU) -> dict[str, Any]:
    """Dispatchable COLLECT: when the directory only exists on another host
    (e.g. ``/Volumes/harqis-data`` from a Windows box), send this to the HFL
    queue so it runs on harqis-server, and read the returned manifest back.

    Returns the full manifest dict (so a ``.get()`` caller gets the material);
    the digest + artifacts are written on the executing host too.
    """
    window = parse_period(period)
    manifest = collect_archive(
        root, window,
        max_files=max_files, max_caption_files=max_caption_files,
        frames_per_video=frames_per_video, do_caption=do_caption,
        caption_model=model, cfg_id=cfg_id__anthropic,
    )
    _write_artifacts(manifest, render_digest(manifest), _slug(window.label))
    return manifest
